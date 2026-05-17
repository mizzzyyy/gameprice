"""Собрать GamePrice_Watcher.pptx из тех же слайдов. Запуск: python build_pptx.py"""
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("Установи: pip install python-pptx")
    raise

OUT = Path(__file__).parent / "GamePrice_Watcher.pptx"
BG = RGBColor(0x0D, 0x11, 0x17)
ACCENT = RGBColor(0x58, 0xA6, 0xFF)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x8B, 0x94, 0x9E)

SLIDES = [
    (
        "GamePrice Watcher",
        [
            "Telegram-бот для мониторинга цен на игры",
            "",
            "Выполнил: [ФИО]",
            "Группа: [№]",
            "Руководитель: [ФИО]",
            "[Год]",
        ],
    ),
    (
        "Актуальность и проблема",
        [
            "• Цены разбросаны: Steam, ключи, RU-магазины",
            "• Разные валюты и региональные ограничения",
            "• Нет единого инструмента в мессенджере",
            "",
            "Идея: один бот — поиск, сравнение, уведомления",
        ],
    ),
    (
        "Цель и задачи",
        [
            "Цель: агрегировать цены и уведомлять о снижении",
            "",
            "• Поиск по названию / URL Steam",
            "• Сбор цен из API и сайтов",
            "• Wishlist с порогом в ₽",
            "• История и график",
            "• RU / EN интерфейс",
        ],
    ),
    (
        "Технологический стек",
        [
            "Python 3 · aiogram 3 · aiohttp",
            "SQLite (aiosqlite) · APScheduler",
            "matplotlib · python-dotenv",
        ],
    ),
    (
        "Архитектура",
        [
            "Пользователь → bot.py → parser.py → API",
            "              ↓",
            "        db.py (SQLite) ← scheduler.py",
            "",
            "Модули: ru_retailers, currency, i18n, config",
        ],
    ),
    (
        "Источники данных",
        [
            "• CheapShark — ключи (USD → ₽, курс ЦБ)",
            "• Steam Store API — RU / US",
            "• Plati — XML API",
            "• Kupikod, Steambuy, GGSEL — HTML",
            "",
            "Минимум = min по всем источникам в ₽",
        ],
    ),
    (
        "Поиск (/search)",
        [
            "• Алиасы + fuzzy-поиск",
            "• Параллельный сбор цен",
            "• Сортировка от дешёвого",
            "• DLC, AppID, CheapShark",
        ],
    ),
    (
        "Wishlist и планировщик",
        [
            "/add — URL + порог в рублях",
            "",
            "Wishlist — каждые 2 ч",
            "Раздачи — 24 ч",
            "Патчи Steam — 6 ч",
            "",
            "цена ≤ порог → уведомление",
        ],
    ),
    (
        "База данных",
        [
            "users · tracked_games · price_history",
            "favorites · freebie_cache",
            "",
            "Файл gameprice.db",
            "Параметризованные SQL-запросы",
        ],
    ),
    (
        "Сложности и решения",
        [
            "Steam из РФ → HTTP_PROXY",
            "Мусор в парсинге → фильтры",
            "Неверный минимум → единый список в ₽",
            "Автоссылка Plati → разрыв автолинка",
        ],
    ),
    (
        "Результаты",
        [
            "✓ Рабочий Telegram-бот",
            "✓ 4+ источника цен",
            "✓ Автоуведомления wishlist",
            "✓ График, 2 языка",
            "",
            "Ограничения: SQLite, парсинг HTML",
        ],
    ),
    (
        "Спасибо за внимание!",
        [
            "Вопросы?",
            "",
            "Демо: /search · /add · /list",
            "Перспективы: VPS, PostgreSQL, Docker",
        ],
    ),
]


def add_slide(prs: Presentation, title: str, lines: list[str]) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(9), Inches(5))
    tf2 = body.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(lines):
        para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        para.text = line
        para.font.size = Pt(20 if line.startswith("•") or line.startswith("✓") else 22)
        para.font.color.rgb = MUTED if line == "" else TEXT
        if "Цель" in line or "Идея" in line:
            para.font.bold = True
            para.font.color.rgb = ACCENT


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    for title, lines in SLIDES:
        add_slide(prs, title, lines)
    prs.save(OUT)
    print(f"Готово: {OUT}")


if __name__ == "__main__":
    main()
